import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import re
from groq import Groq

# Establece la clave API para acceder a la API de Groq desde st.secrets
api_key = st.secrets["general"]["GROQ_API_KEY"]


st.title("Generador de Presentaciones")
tema_input = st.text_input("Tema: Introduce el tema sobre el que quieres desarrollar la ppt", "Historia del Arte")
cantidad_slides_input = st.selectbox("Cantidad de Slides: Selecciona la cantidad de Hojas que quieres que tenga la ppt", ["5", "2", "3", "4", "6", "7", "8", "9"])
publico_objetivo_input = st.text_input("Público Objetivo: ¿A quien irá dirigida?", "Público en General")
fuente_input = st.text_input("Fuentes de preferencia: Ingresa la fuente de preferencia, por ejemplo: publicaciones de organización x", "Lo que encuentres")


# Barra lateral
with st.sidebar:
    st.write("Estás usando **Streamlit💻** and **Groq🖥**\n from Vitto ✳️")
    
    uploaded_file = st.file_uploader("Si subes un txt la ppt se genera con estos datos", type=["txt"])

    modelo = st.selectbox("Modelo", ["llama-3.1-70b-versatile", "llama3-70b-8192", "mixtral-8x7b-32768"])

    max_tokens = st.selectbox("Max New Tokens", [4096, 2048, 1024])  

    temperature = st.slider("Temperatura", 0.0, 1.0, 0.5, 0.2)


def llama3(prompt, modelo:str="llama-3.1-70b-versatile", max_tokens:int=4096, temperature:int=0.5):
    client = Groq(api_key = api_key)
    MODEL = modelo
    # Step 1: send the conversation and available functions to the model
    messages=[
        {
            "role": "system",
            "content": "Eres un Asistente experto"
        },
        {
            "role": "user",
            "content": prompt,
        }
    ]
    
    response = client.chat.completions.create(
        model=MODEL,
        messages=messages,
        #tools=tools,
        temperature=temperature,
        tool_choice="auto",
        max_tokens=max_tokens
    )

    response_message = response.choices[0].message.content
    
    return response_message

def analizar_fuente():
    
    if uploaded_file is not None:
        # Leer el contenido del archivo .txt
        fuente = uploaded_file.read().decode("utf-8")
        nombre_archivo = uploaded_file.name
        st.write(f"La fuente cargada es: {nombre_archivo}")
        #st.write(f"La fuente cargada es: {fuente}")
    else:
        fuente = fuente_input
    
    return fuente

def update_progress_bar(value):
    progress_bar.progress(value)

def eliminar_asteriscos(texto):
    # Utilizar una expresión regular para encontrar asteriscos
    asterisco_pattern = re.compile(r'\*', flags=re.UNICODE)
    # Eliminar asteriscos del texto
    texto_sin_asteriscos = asterisco_pattern.sub('', texto)
    return texto_sin_asteriscos    

def generar_presentacion():
    update_progress_bar(5)
     

    # Obtener los valores de las cajas de texto en Streamlit
    tema = tema_input
    cantidad_slides = int(cantidad_slides_input)
    publico_objetivo = publico_objetivo_input
    fuentes = analizar_fuente()
    update_progress_bar(25)

    prompt = f""" 
    Genera mucho contenido en español-argentina para una presentación en PowerPoint (aproximadamente 300 palabras por slide).
     *Enumera los slides como Slide 1, Slide 2, Slide 3, etc.
     *Indica el título de cada slide con el formato "Título:", sin comillas.
     *Crea títulos breves, impactantes y claros (menos de 8 palabras).
     *La estructura narrativa debe estar basada en: Introducción o planteamiento, Desarrollo o nudo, Clímax y Desenlace. 
     *La introducción debe ser motivadora y mencionar: "A continuación expongo hallazgos en la voz de nuestros colaboradores de la 
      última semana". Y luego un breve resumen de lo que se expondrá en toda la presentación, con algún gatillo al desarrollo y 
      Climax, que generen espectativas, ganas de continuar leyendo. (Usa tres Bloques de texto con saltos de linea entre ellos)
     *Cada Slide debe tener un párrafo de introducción a la información que presentes, relacionado con el título y los datos del slide.
     *Evita referenciar mis peticiones literales de "Introducción:", "Desarrollo:", "Desenlace:", "Clímax:", "Clímax y Desenlaze:", "Resumen:", "Introduccion Motivadora:".
     *Otorga dinamismo a la estructura del texto, el parrafo introductorio en un bloque por ejemplo, y si listas datos 
      que tengan una sangria más pronunciada usando viñetas eventualmente. Se creativa!
     *No abuses de un recurso, si usas análisis estadísticos por ejemplo, porcentajes, sumas, promedios, tendencias. No lo hagas en más
      de dos slides de la misma forma. 
     *En el ultimo slide -el que expones las conclusiones- anexa al contenido alguna frase reconocida que tenga relación analógica o metafórica con el informe, cita su autor.  

    Detalles de la solicitud
     *Tema: {tema}
     *Cantidad de Slides: {cantidad_slides}
     *Público objetivo: {publico_objetivo}
     *Fuentes: {fuentes}
    """

    # Generar contenido
    prompt = prompt
    contenido = llama3(prompt, modelo, max_tokens, temperature)
    contenidosa = eliminar_asteriscos(contenido)
    update_progress_bar(50)

    # Separar el contenido en diapositivas
    slides = contenidosa.split("Slide")

    # Crear una presentación PowerPoint
    prs = Presentation()
    
    # Configurar el tamaño de la diapositiva a 16:9 (25.4 cm x 14.29 cm)
    prs.slide_width = Inches(13.33)  # 25.4 cm
    prs.slide_height = Inches(7.5)   # 14.29 cm

    update_progress_bar(75)

    # Agregar cada diapositiva a la presentación

    for slide_text in slides[1:]:  # El primer elemento de la lista está vacío
        titulo = slide_text.split("Título:", 1)
        title = titulo[1].split("\n", 1)[0].strip()
        content = titulo[1].split("\n", 1)[1].strip()
        
        # Agregar una diapositiva
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Puedes elegir un diseño de diapositiva diferente si lo deseas
    
        # Configurar el título
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = 'Poppins'
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
        # Configurar el contenido
        content_shape = slide.placeholders[1]
        content_shape.text = content
    
        # Reducir el tamaño de la fuente del contenido
        for paragraph in content_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)

    # Guardar la presentación
    prs.save('presentacion_generada.pptx')
    update_progress_bar(100)


if st.button("Generar Presentación"):
    progress_bar = st.progress(0)
    
    try:
        # Llamada a la función para generar la presentación
        generar_presentacion()
        
        # Actualiza la barra de progreso al 100%
        progress_bar.progress(100)
        
        # Mensaje de éxito
        st.success("Presentación generada exitosamente. Descárgala desde abajo.")
        
        # Opción para descargar la presentación
        with open("presentacion_generada.pptx", "rb") as file:
            st.download_button("Descargar Presentación", file, "presentacion_generada.pptx")
    
    except Exception as e:
        # Mensaje de error
        st.error(f"Se produjo un error al generar la presentación: {e}\n\nIntenta de nuevo o selecciona un modelo diferente.")
